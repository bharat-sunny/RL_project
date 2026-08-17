"""Check the generated deck for text that runs off the slide or collides.

The deck is laid out with absolute coordinates and an estimate of how tall each
wrapped text block will be.  That estimate is the fragile part: if it is too
small, blocks overlap; if too large, the last block falls off the bottom.  This
script re-derives the rendered height of every text box from its width, font
size and character count, and reports anything that overflows the slide or
overlaps its neighbour.

    python slides/check_layout.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

DECK = Path(__file__).resolve().parent / "RL_Project_Part2_Presentation.pptx"

# Helvetica-like proportional text averages a little under half the point size
# per character; 0.50 is deliberately pessimistic so the check errs toward
# flagging rather than missing.
CHAR_WIDTH_RATIO = 0.50
LINE_HEIGHT_RATIO = 1.22
BOTTOM_MARGIN_IN = 0.25


def estimate_height_emu(text: str, font_pt: float, width_emu: int,
                        line_spacing: float) -> int:
    width_pt = width_emu / 12700.0
    chars_per_line = max(1, int(width_pt / (font_pt * CHAR_WIDTH_RATIO)))

    lines = 0
    for paragraph in text.split("\n"):
        lines += max(1, -(-len(paragraph) // chars_per_line))  # ceil division
    return int(lines * font_pt * LINE_HEIGHT_RATIO * line_spacing * 12700)


def main() -> int:
    if not DECK.exists():
        print(f"no deck at {DECK} — run python slides/build_slides.py first")
        return 1

    prs = Presentation(DECK)
    slide_h = prs.slide_height
    slide_w = prs.slide_width
    limit = slide_h - Emu(int(BOTTOM_MARGIN_IN * 914400))

    problems = 0
    for index, slide in enumerate(prs.slides, start=1):
        boxes = []
        for shape in slide.shapes:
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
            if not runs:
                continue
            font_pt = max((r.font.size.pt for r in runs if r.font.size), default=12)
            spacing = max((p.line_spacing or 1.0)
                          for p in shape.text_frame.paragraphs)
            height = estimate_height_emu(shape.text_frame.text, font_pt,
                                         shape.width, spacing)
            boxes.append({
                "text": shape.text_frame.text[:52].replace("\n", " "),
                "top": shape.top, "bottom": shape.top + height,
                "left": shape.left, "right": shape.left + shape.width,
                "size": font_pt,
            })

        for shape in slide.shapes:
            if shape.shape_type == 13:  # picture
                if shape.top + shape.height > limit:
                    print(f"  slide {index:>2}: IMAGE overflows bottom by "
                          f"{(shape.top + shape.height - slide_h) / 914400:.2f} in")
                    problems += 1
                if shape.left < 0 or shape.left + shape.width > slide_w:
                    print(f"  slide {index:>2}: IMAGE outside horizontal bounds")
                    problems += 1

        for box in boxes:
            if box["bottom"] > limit:
                print(f"  slide {index:>2}: TEXT overflows bottom by "
                      f"{(box['bottom'] - slide_h) / 914400:.2f} in "
                      f"({box['size']:.0f}pt) — “{box['text']}…”")
                problems += 1
            if box["right"] > slide_w:
                print(f"  slide {index:>2}: TEXT past right edge — “{box['text']}…”")
                problems += 1

        ordered = sorted(boxes, key=lambda b: b["top"])
        for a, b in zip(ordered, ordered[1:]):
            horizontal_overlap = a["left"] < b["right"] and b["left"] < a["right"]
            if horizontal_overlap and b["top"] < a["bottom"] - Emu(20000):
                print(f"  slide {index:>2}: OVERLAP {(a['bottom'] - b['top']) / 914400:.2f} in "
                      f"— “{a['text']}…” into “{b['text']}…”")
                problems += 1

    print(f"\n{len(prs.slides.__iter__.__self__._sldIdLst)} slides checked, "
          f"{problems} potential layout problem(s)")
    return 1 if problems else 0


if __name__ == "__main__":
    sys.exit(main())
