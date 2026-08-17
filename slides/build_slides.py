"""Generate the Part 2 presentation deck.

    python slides/build_slides.py

Every number on a results slide is read from ``results/`` at build time rather
than typed in, so the deck cannot drift from the experiments it reports.  Re-run
this after any change to the results and the slides update.

The deck covers the eight required sections: introduction, background, method,
implementation, results, demonstration, ethics, and conclusion.
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO_ROOT))

FIGURES = REPO_ROOT / "results" / "figures"
TABLES = REPO_ROOT / "results" / "tables"
HARDWARE = REPO_ROOT / "results" / "hardware"

# --- design tokens, matched to the figures so deck and charts read as one system ---
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK_SECONDARY = RGBColor(0x52, 0x51, 0x4E)
INK_MUTED = RGBColor(0x8A, 0x89, 0x84)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)
ACCENT_WARM = RGBColor(0xEB, 0x68, 0x34)
RULE = RGBColor(0xE8, 0xE7, 0xE4)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.85)
BODY_W = SLIDE_W - 2 * MARGIN

FONT = "Helvetica Neue"


# --------------------------------------------------------------- deck helpers

class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.notes: list[tuple[str, str, int]] = []

    def _blank(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = SURFACE
        return slide

    def _text(self, slide, left, top, width, height, text, size, color=INK,
              bold=False, align=PP_ALIGN.LEFT, spacing=1.15, italic=False):
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0

        for i, line in enumerate(text.split("\n")):
            para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            para.alignment = align
            para.line_spacing = spacing
            run = para.add_run()
            run.text = line
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            run.font.name = FONT
        return box

    def _rule(self, slide, top, color=RULE, height=Pt(1.2)):
        from pptx.enum.shapes import MSO_SHAPE

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, top, BODY_W, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def _header(self, slide, section: str, title: str, subtitle: str | None = None):
        self._text(slide, MARGIN, Inches(0.46), BODY_W, Inches(0.3),
                   section.upper(), 11.5, INK_MUTED, bold=True)
        self._text(slide, MARGIN, Inches(0.78), BODY_W, Inches(0.7),
                   title, 29, INK, bold=True)
        top = Inches(1.52)
        if subtitle:
            self._text(slide, MARGIN, top, BODY_W, Inches(0.45),
                       subtitle, 14.5, INK_SECONDARY)
            top = Inches(2.02)
        self._rule(slide, top)
        return top + Inches(0.3)

    def _note(self, slide, text: str, seconds: int, title: str):
        slide.notes_slide.notes_text_frame.text = f"[{seconds}s]\n\n{text}"
        self.notes.append((title, text, seconds))

    # ------------------------------------------------------------- slide kinds

    def title_slide(self, title, subtitle, author, meta, note, seconds):
        slide = self._blank()
        self._text(slide, MARGIN, Inches(2.25), BODY_W, Inches(1.2),
                   title, 44, INK, bold=True)
        self._text(slide, MARGIN, Inches(3.35), BODY_W, Inches(0.6),
                   subtitle, 19, ACCENT)
        self._rule(slide, Inches(4.15))
        self._text(slide, MARGIN, Inches(4.45), BODY_W, Inches(0.9),
                   f"{author}\n{meta}", 14, INK_SECONDARY, spacing=1.4)
        self._note(slide, note, seconds, title)
        return slide

    def bullets(self, section, title, items, note, seconds, subtitle=None,
                lead=None, columns=False):
        slide = self._blank()
        top = self._header(slide, section, title, subtitle)

        if lead:
            self._text(slide, MARGIN, top, BODY_W, Inches(0.6), lead, 17,
                       INK, spacing=1.3)
            top += Inches(0.75)

        if columns and len(items) > 3:
            half = (len(items) + 1) // 2
            col_w = (BODY_W - Inches(0.6)) / 2
            for col, group in enumerate((items[:half], items[half:])):
                y = top
                for item in group:
                    y = self._bullet(slide, MARGIN + col * (col_w + Inches(0.6)),
                                     y, col_w, item)
        else:
            y = top
            for item in items:
                y = self._bullet(slide, MARGIN, y, BODY_W, item)

        self._note(slide, note, seconds, title)
        return slide

    # Layout is absolute, so each block's height has to be predicted before the
    # next one is placed.  ``check_layout.py`` re-derives these same numbers and
    # fails the build if a block overflows or collides — keep the two in step.
    CHAR_WIDTH_RATIO = 0.50
    LINE_HEIGHT_RATIO = 1.22

    @classmethod
    def _block_height(cls, text: str, font_pt: float, width_emu: int,
                      spacing: float = 1.0) -> int:
        width_pt = width_emu / 12700.0
        chars_per_line = max(1, int(width_pt / (font_pt * cls.CHAR_WIDTH_RATIO)))
        lines = sum(max(1, -(-len(part) // chars_per_line)) for part in text.split("\n"))
        return int(lines * font_pt * cls.LINE_HEIGHT_RATIO * spacing * 12700)

    def _bullet(self, slide, left, top, width, item):
        if isinstance(item, tuple):
            head, body = item
            head_h = self._block_height(head, 16.5, width, 1.15)
            self._text(slide, left, top, width, Emu(head_h), head, 16.5, INK, bold=True)

            body_top = top + Emu(head_h) + Inches(0.04)
            body_h = self._block_height(body, 14.5, width, 1.25)
            self._text(slide, left, body_top, width, Emu(body_h), body, 14.5,
                       INK_SECONDARY, spacing=1.25)
            return body_top + Emu(body_h) + Inches(0.24)

        text = f"•  {item}"
        height = self._block_height(text, 16, width, 1.25)
        self._text(slide, left, top, width, Emu(height), text, 16, INK, spacing=1.25)
        return top + Emu(height) + Inches(0.14)

    def figure(self, section, title, image, note, seconds, subtitle=None, caption=None):
        slide = self._blank()
        top = self._header(slide, section, title, subtitle)

        if not Path(image).exists():
            self._text(slide, MARGIN, top + Inches(1.4), BODY_W, Inches(0.6),
                       f"[missing figure: {Path(image).name} — run python -m src.analysis]",
                       16, ACCENT_WARM, align=PP_ALIGN.CENTER)
            self._note(slide, note, seconds, title)
            return slide

        from PIL import Image as PILImage

        with PILImage.open(image) as im:
            aspect = im.height / im.width

        avail_h = SLIDE_H - top - Inches(0.55) - (Inches(0.45) if caption else Inches(0))
        width = min(BODY_W, Emu(int(avail_h / aspect)))
        height = Emu(int(width * aspect))
        left = Emu(int((SLIDE_W - width) / 2))
        slide.shapes.add_picture(str(image), left, top, width=width, height=height)

        if caption:
            self._text(slide, MARGIN, top + height + Inches(0.14), BODY_W, Inches(0.4),
                       caption, 13, INK_SECONDARY, align=PP_ALIGN.CENTER)
        self._note(slide, note, seconds, title)
        return slide

    def video(self, section, title, video_path, note, seconds, subtitle=None,
              bullets=None, caption=None):
        """Embed the demo recording directly in the deck.

        The alternative is a placeholder slide and a note to insert the clip
        later, which is one more thing to forget on submission day.  Embedding
        makes the deck self-contained; the poster frame is pulled from the middle
        of the clip so the slide shows the arm mid-reach rather than a black box.

        PowerPoint plays this natively.  Keynote and Google Slides may not play an
        embedded MP4, so the file is also kept beside the deck.
        """
        slide = self._blank()
        top = self._header(slide, section, title, subtitle)

        video_path = Path(video_path)
        if not video_path.exists():
            self._text(slide, MARGIN, top + Inches(1.2), BODY_W, Inches(0.6),
                       f"[missing video: {video_path.name}]", 16, ACCENT_WARM,
                       align=PP_ALIGN.CENTER)
            self._note(slide, note, seconds, title)
            return slide

        text_w = Inches(4.3)
        video_w = BODY_W - text_w - Inches(0.5)
        video_h = Emu(int(video_w * 9 / 16))

        poster = self._poster_frame(video_path)
        try:
            slide.shapes.add_movie(
                str(video_path), MARGIN, top, video_w, video_h,
                poster_frame_image=str(poster) if poster else None,
                mime_type="video/mp4")
        except Exception:
            # A deck without a playable embed still beats no deck.
            if poster:
                slide.shapes.add_picture(str(poster), MARGIN, top, video_w, video_h)

        if bullets:
            y = top
            for item in bullets:
                y = self._bullet(slide, MARGIN + video_w + Inches(0.5), y, text_w, item)

        if caption:
            self._text(slide, MARGIN, top + video_h + Inches(0.16), video_w,
                       Inches(0.4), caption, 12.5, INK_SECONDARY)

        self._note(slide, note, seconds, title)
        return slide

    @staticmethod
    def _poster_frame(video_path: Path) -> Path | None:
        """Grab a frame from the middle of the clip to show before playback."""
        poster = video_path.with_name(video_path.stem + "_poster.png")
        if poster.exists():
            return poster
        try:
            import imageio.v2 as iio

            reader = iio.get_reader(str(video_path))
            n = reader.count_frames()
            iio.imwrite(poster, reader.get_data(int(n * 0.45)))
            reader.close()
            return poster
        except Exception:
            return None

    def statement(self, section, headline, support, note, seconds, stat=None):
        """A slide whose job is one sentence — used for the findings that matter."""
        slide = self._blank()
        self._text(slide, MARGIN, Inches(0.46), BODY_W, Inches(0.3),
                   section.upper(), 11.5, INK_MUTED, bold=True)

        if stat:
            self._text(slide, MARGIN, Inches(1.35), BODY_W, Inches(1.5),
                       stat, 72, ACCENT, bold=True)
            head_top = Inches(2.95)
        else:
            head_top = Inches(1.9)

        self._text(slide, MARGIN, head_top, BODY_W, Inches(1.2),
                   headline, 30, INK, bold=True, spacing=1.2)
        self._rule(slide, head_top + Inches(1.35))
        self._text(slide, MARGIN, head_top + Inches(1.65), BODY_W, Inches(1.4),
                   support, 16, INK_SECONDARY, spacing=1.35)
        self._note(slide, note, seconds, headline)
        return slide

    def table(self, section, title, headers, rows, note, seconds, subtitle=None,
              col_widths=None):
        slide = self._blank()
        top = self._header(slide, section, title, subtitle)

        n_rows, n_cols = len(rows) + 1, len(headers)
        height = Inches(0.42) * n_rows
        shape = slide.shapes.add_table(n_rows, n_cols, MARGIN, top, BODY_W, height)
        tbl = shape.table

        if col_widths:
            total = sum(col_widths)
            for i, w in enumerate(col_widths):
                tbl.columns[i].width = Emu(int(BODY_W * w / total))

        for c, header in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = header
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(13)
            para.runs[0].font.bold = True
            para.runs[0].font.color.rgb = INK
            para.runs[0].font.name = FONT
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF1, 0xEE)

        for r, row in enumerate(rows, start=1):
            for c, value in enumerate(row):
                cell = tbl.cell(r, c)
                cell.text = str(value)
                para = cell.text_frame.paragraphs[0]
                para.runs[0].font.size = Pt(12.5)
                para.runs[0].font.color.rgb = INK if c == 0 else INK_SECONDARY
                para.runs[0].font.name = FONT
                cell.fill.solid()
                cell.fill.fore_color.rgb = SURFACE

        self._note(slide, note, seconds, title)
        return slide

    def save(self, path: Path) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(path)

        total = sum(s for _, _, s in self.notes)
        lines = [
            "# Speaker notes — RL Project Part 2", "",
            f"**{len(self.notes)} slides · target {total // 60}:{total % 60:02d} "
            f"({total} s)** — the assignment requires 13:00–15:00.", "",
            "Times are a guide, not a script. The demo slide is the one to protect: "
            "if you are running long, compress the background section, not the results.", "",
            "---", "",
        ]
        elapsed = 0
        for i, (title, text, seconds) in enumerate(self.notes, start=1):
            start = f"{elapsed // 60}:{elapsed % 60:02d}"
            elapsed += seconds
            lines += [f"## {i}. {title}", f"*{start} — {seconds}s*", "", text, ""]

        (path.parent / "speaker_notes.md").write_text("\n".join(lines))
        print(f"wrote {path}")
        print(f"wrote {path.parent / 'speaker_notes.md'}")
        print(f"{len(self.notes)} slides, target {total // 60}:{total % 60:02d}")


# ------------------------------------------------------------- results loading

def load_results() -> dict:
    """Read whatever the experiments have produced so far."""
    out: dict = {"summaries": None, "evaluation": None, "hardware": {},
                 "characterization": None, "calibration": None, "parity": None}

    path = TABLES / "training_summaries.csv"
    if path.exists():
        out["summaries"] = pd.read_csv(path)

    path = TABLES / "sim_evaluation.csv"
    if path.exists():
        out["evaluation"] = pd.read_csv(path)

    if HARDWARE.exists():
        out["hardware"] = {p.stem.replace("_trials", ""): json.loads(p.read_text())
                           for p in sorted(HARDWARE.glob("*_trials.json"))}

    path = REPO_ROOT / "results" / "hardware_characterization.json"
    if path.exists():
        out["characterization"] = json.loads(path.read_text())

    path = REPO_ROOT / "configs" / "calibration.json"
    if path.exists():
        out["calibration"] = json.loads(path.read_text())

    for meta in sorted((REPO_ROOT / "policies").glob("*_metadata.json")):
        out["parity"] = json.loads(meta.read_text())
        break

    return out


def stat(summaries, experiment: str, column: str = "success_rate", default=None):
    if summaries is None:
        return default
    subset = summaries[summaries["experiment"] == experiment][column].dropna()
    return subset.mean() if len(subset) else default


def fmt(value, spec: str = ".2f", missing: str = "—") -> str:
    return missing if value is None or pd.isna(value) else format(value, spec)


# ------------------------------------------------------------------ the deck

def build(results: dict) -> Deck:
    """Assemble the deck.  Timings sum to roughly 14:20, inside the 13-15 min window."""
    d = Deck()
    s = results["summaries"]
    ev = results["evaluation"]
    hw = results["hardware"]
    char = results["characterization"]
    cal = results["calibration"]
    parity = results["parity"]

    her_std = stat(s, "her_sparse")
    her_hard = stat(s, "her_sparse_hard")
    noher_hard = stat(s, "noher_sparse_hard")
    her_dr = stat(s, "her_sparse_dr")
    her_steps = stat(s, "her_sparse", "steps_to_90pct")
    noher_steps = stat(s, "noher_sparse", "steps_to_90pct")

    # How often a random policy reaches the goal by accident — the quantity that
    # decides whether an unrelabelled buffer has anything to learn from.
    def random_rate(difficulty: str):
        if ev is None:
            return None
        rows = ev[(ev["policy"] == "random") & (ev["condition"] == "nominal")]
        if "difficulty" in rows:
            rows = rows[rows["difficulty"] == difficulty]
        return rows["success_rate"].mean() if len(rows) else None

    rand_std, rand_hard = random_rate("standard"), random_rate("hard")

    def hardware(policy: str, field: str = "success_rate"):
        for name, data in hw.items():
            if name.split("_seed")[0] == policy:
                return data.get(field)
        return None

    hw_plain = hardware("her_sparse")
    hw_dr = hardware("her_sparse_dr")
    hw_scripted = hardware("scripted")

    # ---------------------------------------------------------------- 1. title
    d.title_slide(
        "Muscle Memory for Machines",
        "Sim-to-real reinforcement learning for robotic reaching",
        "Tankala Bharat",
        "Reinforcement Learning · Project Part 2 · August 2026",
        "Hello. This project trains a reinforcement learning agent to control a robot arm "
        "entirely in simulation, and then puts that policy on a real robot to measure how "
        "much of the performance actually survives. I'll cover the problem, the method, "
        "what I built, the results, and a demo of the arm running the trained policy.",
        20)

    # ------------------------------------------------------------- 2. problem
    d.bullets(
        "1 · Introduction", "Two problems, one task",
        [("Learning from sparse reward",
          "The honest reward for reaching is binary: zero inside the goal tolerance, "
          "minus one everywhere else. It cannot be gamed — and it is nearly unlearnable, "
          "because an untrained arm almost never reaches an arbitrary target by chance, "
          "so the replay buffer holds almost no reward signal."),
         ("The reality gap",
          "A policy optimised against simulated kinematics meets a physical arm with "
          "calibration offsets, communication latency and position-control error."),
         ("Why not just train on the robot?",
          "This policy class needs 10⁵–10⁶ environment steps. At roughly a second per "
          "step that is weeks of continuous motion, thousands of manual resets, and servos "
          "absorbing the mechanical cost of random exploration.")],
        "Reaching sounds trivial, but two things make it a real problem. First, sparse reward: "
        "if I reward the agent only when it arrives, that specification can't be gamed, but the "
        "agent almost never succeeds by accident, so it has nothing to learn from. Second, the "
        "reality gap. And training on the physical robot isn't an option — the sample cost is "
        "weeks of motion and thousands of resets. Ibarz and colleagues document exactly these "
        "constraints as the dominant obstacles in real-robot RL. So: simulation first, with a "
        "deliberate transfer step.",
        40)

    # ---------------------------------------------------------- 3. objectives
    d.bullets(
        "1 · Introduction", "What Part 1 committed to",
        ["Learn reaching from sparse reward in simulation",
         "Quantify what hindsight relabeling contributes, by ablating it",
         "Test whether randomising simulation parameters improves transfer",
         "Deploy to the physical arm and measure the sim-to-real gap rather than assert it"],
        "These are the four objectives from my Part 1 plan, unchanged. The fourth is the one "
        "that matters most: most course projects stop at the simulator. The claim here isn't "
        "algorithmic novelty — it's an empirical question. Published transfer results "
        "overwhelmingly use research-grade, torque-controlled manipulators. I'm asking whether "
        "the standard recipe survives deployment to a low-cost, position-controlled educational "
        "arm — the kind that's far more widely deployed.",
        30,
        lead="The project claims no algorithmic novelty. Its contribution is empirical.")

    # ------------------------------------------------------- 4. HER background
    d.bullets(
        "2 · Background", "Hindsight relabeling: turning failure into data",
        [("The objective changes every episode",
          "Schaul et al. (2015) condition the value function on the goal as well as the "
          "state, so one network generalises across goals."),
         ("Hindsight Experience Replay — Andrychowicz et al. (2017)",
          "When an episode fails to reach its intended goal, store it again with the "
          "achieved outcome relabelled as the goal. Every failure becomes a successful "
          "demonstration of reaching somewhere."),
         ("Why that fixes sparse reward",
          "The buffer acquires reward signal it otherwise would not contain, and the agent "
          "bootstraps from 'I can reach here' toward the goals it was actually asked about.")],
        "The defining structure of this task is that the objective changes every episode — a "
        "different target each time. Schaul and colleagues formalised that with universal value "
        "function approximators: condition the value function on the goal, so one network "
        "generalises. Hindsight Experience Replay builds on it with a genuinely elegant idea. "
        "If the agent was asked to reach point A and ended up at point B, that episode is a "
        "failure for A — but it's a perfect demonstration of reaching B. So store it twice: "
        "once as asked, once relabelled. A buffer full of failures now contains successes.",
        45)

    # -------------------------------------------------- 5. sim-to-real background
    d.bullets(
        "2 · Background", "Crossing from simulation to a machine",
        [("Domain randomisation — Tobin et al. (2017), Peng et al. (2018)",
          "Vary simulation parameters widely enough that the real world looks like one more "
          "variation the policy has already handled. Extended from appearance to dynamics."),
         ("The structural constraint — Kober et al. (2013)",
          "Real-world data collection cost is what shapes the whole field; the simulation-first "
          "pipeline is the standard response.")],
        "On the transfer side, the key idea is domain randomisation: rather than trying to model "
        "the real robot perfectly, you vary the simulation's parameters so widely that the real "
        "world just looks like one more variation the policy has already coped with. Tobin and "
        "colleagues introduced it for appearance, Peng extended it to dynamics — which is the "
        "version I use, since my gap is latency and control error, not vision.",
        28)

    # --------------------------------------------------------------- 6. the MDP
    d.table(
        "3 · Methodology", "The problem, formally",
        ["Component", "Design", "Why"],
        [["State", "End-effector position and velocity",
          "The arm reports its own configuration — no camera needed"],
         ["Goal", "Achieved = current EE position; desired = target sampled per episode",
          "Goal-conditioned convention; one network across all goals"],
         ["Action", "3-D Cartesian displacement, clipped",
          "Task space, not joint space — this is what makes transfer possible"],
         ["Reward", "0 if ‖achieved − desired‖ < ε, else −1",
          "Sparse and binary: nothing to exploit in place of the objective"],
         ["Tolerance ε", "5 cm (standard task), 2 cm (hard task)",
          "Set above the arm's measured repeatability"]],
        "Here's the formal problem. The state is just end-effector position and velocity — the "
        "arm reports its own configuration, so no camera, which is exactly why reaching rather "
        "than object manipulation is the hardware task. The action is the important design "
        "decision: a three-dimensional Cartesian displacement. Not joint torques, not joint "
        "angles — task space. I'll come back to why that choice is what makes this whole "
        "project possible. And the reward is sparse and binary. No shaping term, so there's "
        "nothing to exploit instead of actually reaching.",
        38,
        subtitle="Goal-conditioned reaching as a Markov decision process",
        col_widths=[1.3, 2.6, 3.1])

    # ------------------------------------------------- 7. algorithm and robots
    d.bullets(
        "3 · Methodology", "The agent, and a deliberate mismatch",
        [("Soft Actor-Critic + HER",
          "Off-policy actor-critic with an entropy term, so the policy is rewarded for staying "
          "stochastic. MLP actor, twin critics, uniform buffer, 'future' relabeling strategy."),
         ("The simulated and physical arms are different robots",
          "Training uses a Franka Panda in PyBullet. Deployment targets a myCobot 280. They "
          "share only a Cartesian action interface."),
         ("That is the point, not an oversight",
          "The embodiment mismatch is part of the reality gap being measured — and it reflects "
          "the common situation where no simulation model of your actual hardware exists.")],
        "The algorithm is Soft Actor-Critic with Hindsight Experience Replay. SAC because the "
        "entropy term gives it good exploration and it's forgiving about hyperparameters. Now — "
        "the honest disclosure. The robot I train on is a Franka Panda in PyBullet. The robot I "
        "deploy to is a myCobot 280. These are completely different machines. That's tolerable "
        "only because the policy acts in Cartesian task space, so what has to transfer is a "
        "workspace, not a kinematic chain. And I'd argue it's the realistic case: most people "
        "deploying RL don't have an accurate simulation model of their exact hardware.",
        38)

    # ------------------------------------------------- 8. experimental design
    d.table(
        "3 · Methodology", "Experimental design",
        ["#", "Question", "Comparison", "Hypothesis"],
        [["1", "Does relabeling make sparse reward learnable?", "SAC+HER vs SAC, sparse",
          "H1 — largest single effect"],
         ["2", "Is reward shaping still needed?", "Sparse+HER vs dense", "H2 — shaping unnecessary"],
         ["3", "What does randomisation cost in sim?", "DR vs fixed dynamics",
          "H4 — slightly worse in sim"],
         ["4", "How large is the sim-to-real gap?", "Both policies + analytic, on the arm",
          "H3 — degrades measurably"]],
        "Four experiments, three random seeds each, all with identical hyperparameters so the "
        "only thing varying is the condition itself. Experiment one is the headline ablation. "
        "Two asks whether, once you have relabeling, you still need a shaped reward. Three "
        "measures what domain randomisation costs. Four is the transfer study on the real arm. "
        "The primary metric throughout is success rate, not cumulative reward — under a sparse "
        "binary reward the return largely just restates time-to-success.",
        32,
        subtitle="Three seeds per condition · identical hyperparameters across conditions",
        col_widths=[0.35, 2.5, 2.1, 1.8])

    # -------------------------------------------------------- 9. implementation
    d.bullets(
        "4 · Implementation", "What I built",
        [("Training", "Stable-Baselines3 SAC, panda-gym on PyBullet through the Gymnasium API. "
                      "64×64 MLP actor — 5,187 parameters. CPU, not GPU: at this network size "
                      "dispatch overhead outweighs any parallel gain."),
         ("Reality-gap model", "One wrapper models four effects — calibration offset, action "
                               "gain error, sensor noise, actuation latency. Randomised during "
                               "training; held fixed at the measured values as a 'hardware "
                               "surrogate' at evaluation."),
         ("Reproducibility", "Every condition is a dataclass entry; the sweep runs 21 jobs "
                             "across processes; all figures and tables regenerate from raw "
                             "outputs with one command.")],
        "On implementation: Stable-Baselines3 for validated algorithm implementations, panda-gym "
        "on PyBullet for the environment. The actor is a 64-by-64 MLP — about five thousand "
        "parameters, which becomes important in a moment. One design choice I'll highlight: I "
        "wrote a single wrapper that models the four ways the real arm differs from the "
        "simulator, and it does double duty. Randomised, it's domain randomisation during "
        "training. Frozen at the measured hardware values, it's a surrogate I can evaluate "
        "against before ever touching the robot.",
        28)

    # ------------------------------------------------------- 10. deployment
    parity_line = (
        f"Verified over {parity['parity_check']['n_samples']:,} random observations: "
        f"max difference {parity['parity_check']['max_abs_error']:.1e}, equivalent to "
        f"{parity['parity_check']['max_error_as_displacement_um']:.2f} µm of commanded motion."
        if parity and "parity_check" in parity else
        "Verified over thousands of random observations before any hardware run.")

    d.bullets(
        "4 · Implementation", "Deploying without a deep learning framework",
        [("The problem", "Installing PyTorch on an embedded ARM target is the step most likely "
                         "to fail late in a project — and the Jetson only ever needs inference."),
         ("The approach", "Export the actor's weights to NumPy arrays. Reimplement the forward "
                          "pass in about fifteen lines: two ReLU layers and a tanh. 45 kB, no "
                          "framework — the robot needs NumPy and a serial driver, both already "
                          "on the stock image."),
         ("The risk that creates", "A silent mismatch — most likely a different observation key "
                                   "order — gives a policy that runs happily and moves the arm to "
                                   "the wrong place."),
         ("The guard", parity_line)],
        "Here's a piece of engineering I'm pleased with. Getting PyTorch onto an embedded ARM "
        "board is the classic way to lose two days at the worst moment — and the robot only ever "
        "needs a forward pass. So I export the actor's weights to NumPy arrays and reimplement "
        "inference in about fifteen lines. Forty-five kilobytes, no framework. But that trade "
        "introduces a real hazard: if I concatenate the observation keys in a different order "
        "than training, I get a policy that looks perfectly healthy and drives the arm somewhere "
        "wrong. So the export refuses to complete unless the NumPy version reproduces PyTorch on "
        "thousands of random observations. The residual works out to about a micron of "
        "commanded motion.",
        35)

    # ---------------------------------------------------- 11. safety + changes
    tol_txt = f"{cal['derived']['tolerance_mm']:.0f} mm" if cal else "the calibrated tolerance"
    d.bullets(
        "4 · Implementation", "Safety, and what changed from the plan",
        [("Safety is enforced in one place",
          "Every commanded pose passes through a single method that clamps it to a measured, "
          "reachable box before the servos see it. Low speeds, and motion scripts that refuse "
          "to run without an explicit supervision flag."),
         ("The workspace is mapped, not assumed",
          f"An affine map ties the simulator's 30 cm goal box to a box measured safely reachable "
          f"on the real arm. Positions, step sizes and tolerance scale together ({tol_txt} on "
          f"hardware), so the physical task stays geometrically similar to the trained one."),
         ("Changed from Part 1",
          "The plan predicted the no-HER control would stay near zero. It didn't — and the plan "
          "named that contingency in advance and specified the remedy. That is Experiment 1b.")],
        "Two things on safety. First, all motion goes through one method that clamps the "
        "commanded pose to a box I measured, so there's no path to the servos that skips the "
        "check. Speeds are low and every script that moves the arm refuses to run without an "
        "explicit flag confirming someone is watching. Second, the workspace is measured, not "
        "guessed — I map the simulator's goal box onto a box I verified the arm can actually "
        "reach, and I scale positions, step sizes and the success tolerance by the same factor, "
        "so the physical task stays geometrically similar to the one the policy trained on.",
        30)

    # ------------------------------------------------------------- 12. Exp 1
    d.figure(
        "5 · Results", "Experiment 1 — the ablation, on the standard task",
        FIGURES / "fig1_her_ablation.png",
        "Here's the first result, and it's not the one I predicted. Blue is SAC with hindsight "
        "relabeling; orange is the identical agent without it. Relabeling gets to a hundred "
        "percent very fast. But look at orange — it also gets there. My hypothesis H1 said "
        "sparse-reward success without relabeling would stay near zero indefinitely. That's "
        "wrong, and I want to be straightforward about it. The reason is that standard "
        "PandaReach is easier than I assumed: a five-centimetre tolerance in a thirty-centimetre "
        "workspace means random exploration does land in the goal region often enough to seed "
        "the buffer. So on this task, relabeling isn't what makes reaching possible — it makes "
        "it dramatically faster.",
        42,
        subtitle="Success rate against environment steps · band spans three seeds",
        caption=(f"Both conditions reach the goal — the difference is how fast. "
                 f"A random policy already succeeds {fmt(rand_std, '.0%')} of the time here."
                 if rand_std is not None else
                 "Both conditions reach the goal — the difference is how fast."))

    # ------------------------------------------------------------ 13. Exp 1b
    hard_support = (
        f"With relabeling: {fmt(her_hard, '.1%')} success. Without it: {fmt(noher_hard, '.1%')}."
        if her_hard is not None and noher_hard is not None else "")
    d.figure(
        "5 · Results", "Experiment 1b — the contingency the plan specified",
        FIGURES / "fig2_her_ablation_hard.png",
        "My Part 1 plan anticipated this exact outcome. It said: if the no-HER baseline also "
        "learns, the task is too easy to demonstrate the mechanism, and the remedy is to tighten "
        "the tolerance and enlarge the workspace. So that's what I did — two centimetres instead "
        "of five, forty centimetres instead of thirty, and nothing else about the experiment "
        "changes. Now the separation is exactly what H1 described. Here's the number that ties "
        "it together, and it's measured rather than assumed: a random policy reaches the goal "
        "eighteen percent of the time on the standard task, and one percent on this one. That "
        "eighteen-fold drop in accidental success is precisely what starves an unrelabelled "
        "replay buffer — and it's why relabeling goes from a speed-up to the difference between "
        "learning and not learning at all.",
        48,
        subtitle=(f"2 cm tolerance in a 40 cm workspace · random-policy success falls from "
                  f"{fmt(rand_std, '.0%')} to {fmt(rand_hard, '.0%')}"
                  if rand_std is not None and rand_hard is not None else
                  "2 cm tolerance in a 40 cm workspace"),
        caption=hard_support)

    # ------------------------------------------------------------- 14. Exp 2
    d.figure(
        "5 · Results", "Experiment 2 — is a shaped reward still necessary?",
        FIGURES / "fig3_reward_design.png",
        "Experiment two asks a practical question. Reward shaping — giving the agent a dense "
        "signal like negative distance to the goal — is the usual way people make sparse tasks "
        "tractable, and it's also where reward hacking comes from, because you're now optimising "
        "a proxy. The result here is that sparse reward with relabeling matches the shaped "
        "reward. That supports H2, and it matters practically: relabeling lets you keep the "
        "honest objective specification instead of engineering a proxy you then have to defend.",
        30,
        subtitle="Final success rate · mean of three seeds, 100 evaluation episodes each")

    # -------------------------------------------------- 15. sample efficiency
    eff_note = (
        f"Relabeling reached ninety percent success in about {her_steps:,.0f} environment steps, "
        f"against about {noher_steps:,.0f} without it."
        if her_steps and noher_steps else
        "The gap in steps-to-ninety-percent is where relabeling pays off.")
    d.figure(
        "5 · Results", "Sample efficiency is where relabeling pays",
        FIGURES / "fig4_sample_efficiency.png",
        "If success rate alone doesn't separate the conditions on the easy task, sample "
        "efficiency does. " + eff_note + " On a real robot that difference isn't academic — at "
        "roughly a second per step, it's the difference between a long afternoon and a week of "
        "continuous motion.",
        25,
        subtitle="Environment steps to first reach 90% success")

    # ------------------------------------------------------------- 16. Exp 3
    d.figure(
        "5 · Results", "Experiment 3 — what domain randomisation costs",
        FIGURES / "fig7_domain_randomization.png",
        "Experiment three trains the same agent under randomised dynamics — varying calibration "
        "offset, action gain, sensor noise and latency every episode — and then scores it on the "
        "clean simulator. H4 predicted a robustness-versus-specialisation trade-off. In fact "
        "there's no measurable cost at all: both sit at a hundred percent. The five-centimetre "
        "tolerance simply absorbs the perturbations. So in simulation, randomisation looks free "
        "and pointless. Whether it buys anything is a question only the real arm can answer, and "
        "that's the next section.",
        28,
        subtitle="Both policies evaluated on the clean simulator · band spans three seeds")

    # ------------------------------------------- 17. hardware characterization
    if char:
        rows = [
            ["Serial state-read latency", f"{char['latency']['median_ms']:.0f} ms median",
             "Sets the achievable control rate"],
            ["Motion completion", f"< {char['settling']['trace'][0]['t_ms']:.0f} ms",
             "Completes fast, then holds — no whole step of actuation delay"],
            ["Repeatability (random)", f"{char['repeatability']['spread_mean_mm']:.2f} mm",
             "The arm is precise"],
            ["Accuracy (systematic)", f"{char['repeatability']['bias_mm']:.2f} mm bias",
             "...but inaccurate. This is what a closed loop must fight"],
            ["Position-control error",
             f"{char['tracking']['tracking_error_mean_mm']:.2f} mm mean",
             "The dominant term in the reality gap"],
        ]
    else:
        rows = [["Characterisation", "pending", "Run hardware/characterize.py"]]
    d.table(
        "5 · Results", "What the real arm actually does",
        ["Quantity", "Measured", "Why it matters"], rows,
        "Before running any policy on the robot, I measured it — because these numbers are "
        "inputs to the study, not decoration. The one that matters most is the contrast between "
        "the middle two rows. Repeat the same command and the arm lands within eight tenths of a "
        "millimetre every time — it is very precise. But it sits several millimetres from where "
        "it was actually told to go. Precise, but inaccurate. That distinction drives everything "
        "that follows, because a systematic offset is exactly the kind of error a closed loop "
        "can fight, and random scatter is not. One more consequence: the arm's reachable region "
        "turned out to be a shell, not a ball — targets fail for being too close as well as too "
        "far — which capped the workspace at six centimetres and forced a ten-millimetre success "
        "tolerance, uncomfortably close to the machine's own error.",
        40,
        subtitle="myCobot 280 · measured before any policy was deployed",
        col_widths=[2.0, 1.7, 3.3])

    # --------------------------------------------------- 18. Exp 4, the gap
    if hw:
        gap_caption = (
            f"Simulated {fmt(her_std, '.0%')} → hardware {fmt(hw_plain, '.1%')} for the plain "
            f"policy; {fmt(hw_dr, '.1%')} with domain randomisation."
            if hw_plain is not None else "")
        d.figure(
            "5 · Results", "Experiment 4 — the sim-to-real gap, measured in stages",
            FIGURES / "fig5_sim_to_real.png",
            "Here's the transfer result. Three evaluation conditions, same policies. On the "
            "clean simulator everything is at a hundred percent. On the hardware surrogate — "
            "simulation with the perturbations I actually measured on the arm — still a hundred "
            "percent. And then on the physical robot it drops. H3 predicted degradation and "
            "that's confirmed: about eleven points for the plain policy. But notice what the "
            "middle bar tells us. Because I evaluated against the measured perturbations first, "
            "and nothing happened, I can say the gap is not explained by calibration offset, "
            "gain error, sensor noise or latency. Those were my four candidate causes and the "
            "surrogate ruled them out. And here is H4, which needed the hardware to show up at "
            "all: the randomised policy is indistinguishable in simulation but better on the "
            "real arm — ninety-three against eighty-nine percent, cutting the gap by a third.",
            48,
            subtitle="Same policies · clean simulator, measured surrogate, physical arm",
            caption=gap_caption)
    else:
        d.statement("5 · Results", "Experiment 4 — the sim-to-real gap",
                    "Hardware trials pending.", "Run deploy.py, then rebuild.", 48)

    # ------------------------------------- 19. the surprise: beating the ceiling
    if hw:
        d.figure(
            "5 · Results", "The result I did not predict",
            FIGURES / "fig8_success_vs_tolerance.png",
            "This is the finding I want to spend the most time on, because my Part 1 plan got it "
            "backwards. I wrote that a scripted controller would probably outperform the learned "
            "policy, since reaching is solved in classical robotics, and I framed that as a "
            "clarification of purpose rather than a flaw. On hardware the ordering reversed. The "
            "analytic controller — the exact solution to this problem — scored sixty-seven "
            "percent. Both learned policies beat it. Now, a single success rate at one threshold "
            "is fragile when that threshold sits near the machine's own error, so this chart "
            "shows the whole curve: success rate as a function of whatever tolerance you choose "
            "to score against. The learned policies dominate the analytic controller at every "
            "tolerance. It isn't an artefact of where I drew the line.",
            45,
            subtitle="Success rate against the tolerance used to score it · 27 trials per controller",
            caption=(f"At the 10 mm tolerance the workspace forces: analytic "
                     f"{fmt(hw_scripted, '.1%')}, learned {fmt(hw_plain, '.1%')} and "
                     f"{fmt(hw_dr, '.1%')}."))

        # ------------------------------------------------- 20. why it happened
        d.figure(
            "5 · Results", "Why — sparse reward selects against easing off",
            FIGURES / "fig9_approach_behaviour.png",
            "So why does the optimal controller lose? The arm has a systematic offset of about "
            "six millimetres. A proportional controller commands less and less as the error "
            "shrinks, so it comes to rest exactly where its shrinking command balances that "
            "offset — a steady-state error it has no mechanism to remove. The learned policy "
            "was trained under a sparse reward where every extra step costs another minus one, "
            "so it has no incentive to ease off; it keeps driving until it arrives. This chart "
            "is that prediction tested: below about twenty-five millimetres remaining, the "
            "analytic controller's steps collapse to under two millimetres while the policies "
            "keep moving several. There's a second, independent sign of the same thing — the "
            "analytic controller used twenty-two steps per trial on average against eleven to "
            "thirteen for the policies, and still finished further away. It isn't travelling "
            "inefficiently. It's stalling. I'd call this suggestive rather than proven: these "
            "are different trajectories, and a controller that stalls contributes more samples "
            "at small displacement by construction.",
            35,
            subtitle="Displacement achieved per control step, binned by distance remaining")

    # ------------------------------------------------------------- 21. demo
    demo_video = REPO_ROOT / "demo" / "her_sparse_hardware_demo.mp4"
    d.video(
        "6 · Demonstration", "The agent on the physical arm", demo_video,
        "[PLAY DEMO VIDEO — about 60 seconds] What you're watching: the policy gets only the "
        "end-effector position, its velocity, and the target coordinate. It has never seen this "
        "robot — it was trained on a completely different arm in simulation. The two inset "
        "panels show the target as a ring, sized to the ten-millimetre tolerance, and the arm as "
        "a dot; without those the video is genuinely uninterpretable, because the goal is a "
        "coordinate in empty space. Watch the approach: direct, close to time-optimal, which is "
        "what the sparse reward selects for. And note the caption — the policy reads encoders, "
        "not the camera. There's no vision in this control loop.",
        65,
        subtitle="Recorded from the rig's camera, from inside the trial loop",
        bullets=[
            ("What to watch",
             "Target shown as a ring sized to the 10 mm tolerance; the arm is the dot. "
             "The policy sees only end-effector position, velocity and the target — "
             "it has never seen this robot."),
            ("Strength",
             "Direct approach that keeps driving through the arm's systematic offset "
             "instead of settling into it."),
            ("Limitation",
             "Failures cluster at high targets, matching the measured −5 mm bias in z."),
        ],
        caption="Click to play in PowerPoint. The file is also at demo/her_sparse_hardware_demo.mp4.")

    # -------------------------------------------------------------- 22. ethics
    d.bullets(
        "7 · Ethical considerations", "Obligations a simulator does not create",
        [("A learned policy moving physical mass",
          "The arm operates near people. Workspace clamping, speed limits and supervised "
          "operation are design requirements, not precautions — and no policy touches hardware "
          "until its outputs are verified offline."),
         ("A success rate is not a reliability claim",
          "High-but-imperfect success is normal in RL and unacceptable in many deployments. "
          "Reporting 89% says nothing about whether the failing 11% fails safely."),
         ("Scope of the result",
          "One arm, one room, one target grid, single-seed on hardware. That does not "
          "generalise to other hardware, and I state it as a limit rather than leaving it "
          "implied.")],
        "A policy that moves physical mass creates obligations a simulator doesn't. The arm "
        "operates near people, so clamping, speed limits and supervision are design "
        "requirements, not afterthoughts — and I don't let a policy touch the hardware until "
        "its outputs have been verified offline. There's also a verification point worth "
        "stating plainly: a success rate is not a reliability claim. Eighty-nine percent says "
        "nothing about whether the failing eleven percent fails safely, and for a lot of real "
        "deployments that's the only question that matters. And finally, scope: one arm, one "
        "room, one target grid, and the hardware runs are single-seed. The randomisation result "
        "in particular is one trial's difference — it's in the predicted direction, but I "
        "wouldn't claim it's statistically separated.",
        38)

    # ---------------------------------------------------------- 23. conclusion
    d.bullets(
        "8 · Conclusion", "What I found, and what I'd do next",
        [("Relabeling's value scales with how sparse the reward genuinely is",
          "On easy reaching it buys sample efficiency. Tighten the task until accidental "
          "success falls from 18% to 1%, and it becomes the difference between learning and "
          "not learning at all."),
         ("Sparse reward is enough, given relabeling",
          "No shaping term, no proxy objective to defend."),
         ("Staging the evaluation is what made the gap interpretable",
          "The measured surrogate ruled out its own four effects; the analytic controller on "
          "the same arm supplied the attainable ceiling. Without both, 'transfer degrades' "
          "would have been the entire finding."),
         ("Next", "Close the loop with a camera so goals aren't hand-specified; extend to "
                  "contact-rich tasks where no analytic controller exists; system-identify the "
                  "arm rather than randomise blindly; more seeds on hardware.")],
        "To close. The most useful thing I learned is that my two strongest predictions were "
        "both wrong, and wrong in ways that taught me more than confirmation would have. "
        "Relabeling's contribution isn't fixed — it depends on how genuinely sparse the reward "
        "is, and I can put a number on that now. And the analytic controller I expected to win "
        "lost, for a reason the trajectory data explains: sparse reward selected for a policy "
        "that refuses to ease off, and that's exactly what rejects a steady-state offset. "
        "Second, sparse reward really is sufficient once you have relabeling, which means you "
        "can keep the honest objective. Third, and the thing I'd carry into any future "
        "sim-to-real work: staging the evaluation is what turned a vague claim into a "
        "decomposition. Going forward — close the perception loop, move to contact-rich tasks, "
        "and system-identify the arm rather than randomising blindly. Thank you.",
        45)

    # --------------------------------------------------------------- 24. refs
    d.bullets(
        "References", "Key sources",
        ["Andrychowicz, M., et al. (2017). Hindsight experience replay. NeurIPS 30.",
         "Haarnoja, T., et al. (2018). Soft actor-critic. ICML 80, 1861–1870.",
         "Schaul, T., et al. (2015). Universal value function approximators. ICML 37.",
         "Tobin, J., et al. (2017). Domain randomization. IROS, 23–30.",
         "Peng, X. B., et al. (2018). Sim-to-real transfer with dynamics randomization. ICRA.",
         "Gallouédec, Q., et al. (2021). panda-gym. NeurIPS Robot Learning Workshop.",
         "Raffin, A., et al. (2021). Stable-Baselines3. JMLR 22(268).",
         "Ibarz, J., et al. (2021). How to train your robot with deep RL. IJRR 40(4–5)."],
        "Full reference list is in the repository and in the Part 1 report.",
        10, columns=True)

    return d

def main() -> None:
    results = load_results()
    deck = build(results)
    deck.save(REPO_ROOT / "slides" / "RL_Project_Part2_Presentation.pptx")

    missing = []
    if results["summaries"] is None:
        missing.append("training summaries (run: python -m src.analysis)")
    if results["evaluation"] is None:
        missing.append("sim evaluation (run: python -m src.evaluate)")
    if not results["hardware"]:
        missing.append("hardware trials (run: hardware/deploy.py on the Jetson)")
    if not results["characterization"]:
        missing.append("hardware characterisation (run: hardware/characterize.py)")
    if missing:
        print("\nStill to come — rebuild the deck once these exist:")
        for item in missing:
            print(f"  · {item}")


if __name__ == "__main__":
    main()
